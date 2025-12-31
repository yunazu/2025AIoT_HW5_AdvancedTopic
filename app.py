import streamlit as st
from pptx import Presentation
from google import genai
import io
import json

# --- 介面設定 ---
st.set_page_config(page_title="AI PPT Architect Pro", layout="wide")
st.title("🧠 AI 簡報架構重塑師 (多模板支援版)")

# --- 側邊欄：設定 API 與 說明 ---
with st.sidebar:
    st.header("🔑 設定")
    api_key = st.text_input("輸入 Gemini API Key", type="password")
    st.divider()
    st.markdown("""
    ### 運作原理：
    1. **AI 讀取**：解析舊 PPT 內容。
    2. **AI 重構**：Gemini 重新撰寫精華大綱。
    3. **模板映射**：將新內容注入你上傳的精美模板。
    """)

# --- 核心函式：AI 重組內容 ---
def rewrite_content_with_ai(original_text, api_key):
    client = genai.Client(api_key=api_key)
    
    prompt = f"""
    你是一個專業的簡報架構師。請根據以下原始內容重新設計 3-5 頁簡報大綱。
    
    原始內容：
    {original_text[:4000]}
    
    任務：
    1. 重新梳理內容，精簡為邏輯強、好理解的 3-5 頁大綱。
    2. 嚴格遵守以下 JSON 格式回傳，禁止包含任何 Markdown 標籤或說明文字。
    
    JSON 格式：
    {{
      "slides": [
        {{"title": "標題1", "content": ["重點1", "重點2"]}},
        {{"title": "標題2", "content": ["重點1", "重點2"]}}
      ]
    }}
    """
    response = client.models.generate_content(
            model='gemini-2.5-flash-lite', # Flash 是免費版最穩定的
            contents=prompt
        )
    raw_text = response.text.strip()
    
    # 清理 Markdown 標籤
    if raw_text.startswith("```"):
        raw_text = raw_text.split("```")[1].replace("json", "", 1).strip()
    return json.loads(raw_text)

# --- 核心函式：模板映射與生成 ---
def create_ppt_from_template(ai_data, template_stream):
    # 使用使用者上傳的模板作為基底
    prs = Presentation(template_stream)
    
    for i, slide_data in enumerate(ai_data["slides"]):
        # 選擇版型：通常 0 是標題頁，1 是內容頁
        # 我們假設第一頁用標題頁(0)，其餘用內容頁(1)
        layout_idx = 0 if i == 0 else 1
        try:
            layout = prs.slide_layouts[layout_idx]
        except:
            layout = prs.slide_layouts[0]
            
        slide = prs.slides.add_slide(layout)
        
        # 填入標題
        if slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]
            
        # 填入內容 (尋找內容佔位符)
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1: # 內容區常用的 ID
                shape.text = "\n".join(slide_data["content"])

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# --- 主要 UI 流程 ---
col1, col2 = st.columns(2)

with col1:
    st.subheader("1. 內容來源")
    src_file = st.file_uploader("上傳【原始 PPT】(提取內容用)", type=["pptx"])

with col2:
    st.subheader("2. 視覺風格")
    tpl_file = st.file_uploader("上傳【空白模板】(決定外觀用)", type=["pptx"])

if src_file and tpl_file and api_key:
    if st.button("🚀 開始 AI 重構並更換版型"):
        try:
            with st.spinner("1/2 AI 正在深度閱讀並重構內容..."):
                # 提取舊文字
                old_prs = Presentation(src_file)
                full_text = "\n".join([shape.text for slide in old_prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                
                # 呼叫 AI
                ai_data = rewrite_content_with_ai(full_text, api_key)
            
            with st.spinner("2/2 正在將新內容注入精美模板..."):
                # 生成新 PPT
                result_ppt = create_ppt_from_template(ai_data, tpl_file)
            
            st.success("✅ 簡報重構完成！")
            
            # 預覽內容
            with st.expander("查看 AI 生成的大綱"):
                st.json(ai_data)
                
            st.download_button(
                label="📥 下載重塑後的 PPT",
                data=result_ppt,
                file_name="AI_Pro_Presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            
        except Exception as e:
            st.error(f"發生錯誤：{e}")
            st.info("提示：請確保您的模板中包含標準的『標題』與『內容』佔位符。")

elif not api_key:
    st.warning("👈 請先在側邊欄輸入您的 Gemini API Key。")